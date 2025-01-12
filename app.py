import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from dotenv import load_dotenv
import asyncio

from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai

from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate

# Load environment variables from .env file
load_dotenv()

# Check if GOOGLE_API_KEY is loaded properly
genai_api_key = os.getenv("GOOGLE_API_KEY")
if not genai_api_key:
    raise ValueError("GOOGLE_API_KEY not found in the environment variables.")

# Configure the Google API client
genai.configure(api_key=genai_api_key)

# Function to extract text from PDF files
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Function to extract text from PPTX files
def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Function to split text into chunks (improved logic)
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=113000, chunk_overlap=10000)  # Smaller chunks for better context capture
    chunks = text_splitter.split_text(text)
    return chunks

# Function to create and save vector store (improved embedding model)
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")  # Use a more powerful model here, e.g., text-embedding-ada-002 or gemini
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Function to create conversational chain with more precise prompts
async def get_conversational_chain():
    prompt_template = """
    You are a helpful assistant. Based on the context provided below, answer the user's question in detail. 
    If the answer is not found in the context, say "Answer is not available in the context" and do not make up an answer.

    Context:
    {context}
    
    Question:
    {question}
    
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)  # Ensure you're using the best available model
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user input and generate better responses
async def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")  # You can try different models
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question, k=5)  # Retrieve top 5 similar documents for better context

    # Ensure the response is based on context from the top docs
    chain = await get_conversational_chain()

    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply: ", response["output_text"])

# Main function
def main():
    st.set_page_config("Chat With Multiple PDF and PPTX")
    st.title("Chat with Multiple PDF and PPTX using Gemini")
    st.header("How To Use this App")
    st.write("1. Upload your PDF and PPTX files using the file uploader.")
    st.write("2. Ask a question related to the uploaded files.")
    st.write("3. Click on the 'Submit & Process' button to process the files and get the answer to your question.")

    user_question = st.text_input("Ask a Question from the Files")

    if user_question:
        asyncio.run(user_input(user_question))

    with st.sidebar:
        st.title("Menu:")
    
        uploaded_files = st.file_uploader("Upload your PDF and PPTX Files and Click on the Submit & Process Button", accept_multiple_files=True, type=["pdf", "pptx"])
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                pdf_files = [file for file in uploaded_files if file.type == "application/pdf"]
                pptx_files = [file for file in uploaded_files if file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"]
                
                raw_text = get_pdf_text(pdf_files) + get_pptx_text(pptx_files)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")

    st.write("qwerty")

if __name__ == "__main__":
    main()
